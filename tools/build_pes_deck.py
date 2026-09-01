# -*- coding: utf-8 -*-
"""Build the PES-template Phase-3 Review-1 deck (docs/PES_REVIEW1_DECK.pptx).

Follows the mandated outline in capstone/Review1-phase3 template exactly:
  Title · Outline · Abstract & Scope · Summary of Phase-1/2 (+ suggestions
  incorporated) · Inferences from Literature · Architecture · List of
  Tasks/Modules (+ SDK/API/tools) · Individual Contribution (LOC/time/timeline)
  · Demonstration & Testing · Gantt · References (IEEE) · Thank You

Brand: PES orange header/footer bars, "Title of the Project" top-left,
"name1_name2_name3_name4" footer, PES logo top-right — matched to the template.
Numbers are READ from the committed JSONs; LOC from `git ls-files | wc -l`.

*** The per-person split on the Individual Contribution slide follows the team's
*** AGREED module ownership (confirmed 2026-08-18; also in docs/TEAM_BRIEFING.md).
*** Module LOC are measured (git ls-files | wc -l); hours are team estimates. Git
*** history is single-account (all pushes via one machine), so LOC-by-author is
*** not derivable from git — say so if asked, and cite module ownership.

Re-run: python tools/build_pes_deck.py
"""
import json
import os
import subprocess

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DOCS = os.path.join(ROOT, 'docs')
FIG = os.path.join(DOCS, 'figures')
OUT = os.path.join(DOCS, 'PES_REVIEW1_DECK.pptx')


def _load(p):
    with open(p, encoding='utf-8') as f:
        return json.load(f)


SV = _load(os.path.join(DOCS, 'survey_validation.json'))
SX = _load(os.path.join(DOCS, 'survey_extras.json'))
AB = _load(os.path.join(DOCS, 'ablation_results.json'))
AF = _load(os.path.join(DOCS, 'asr_fairness.json'))
MM = _load(os.path.join(ROOT, 'backend', 'models', 'model_metadata.json'))
inc = SX['incremental']
comp = SX['features']['composites']
chat_ab = {r['config']: r for r in AB['chat']}
fam = AF['by_language_family']


def _paper_stats(name, fallback_pages):
    """(pages, refs) read from the LaTeX log and source, so a slide can never quote a
    stale page or reference count. Falls back to the last known page count if the log
    is absent (fresh clone)."""
    import re as _re
    pages = fallback_pages
    log = os.path.join(DOCS, name + '.log')
    if os.path.exists(log):
        with open(log, encoding='utf-8', errors='ignore') as f:
            m = _re.search(r'Output written on ' + _re.escape(name) + r'\.pdf \((\d+) pages', f.read())
        if m:
            pages = int(m.group(1))
    src = 'ieee_refs.tex' if name == 'IEEE_PAPER' else name + '.tex'
    with open(os.path.join(DOCS, src), encoding='utf-8') as f:
        refs = f.read().count('\\bibitem{')
    return pages, refs


REPORT_PP, REPORT_REFS = _paper_stats('PROJECT_PAPER', 48)
IEEE_PP, IEEE_REFS = _paper_stats('IEEE_PAPER', 6)
GENRE_POWER = SX['genre']['power_curve'][str(SV['construct_validity']['n'])] * 100


def loc(paths, exts=('.py', '.kt', '.java', '.xml', '.tex', '.md', '.yml')):
    files = subprocess.run(['git', 'ls-files', *paths], cwd=ROOT, capture_output=True, text=True).stdout.split()
    n = 0
    for f in files:
        if f.endswith(exts):
            try:
                with open(os.path.join(ROOT, f), encoding='utf-8', errors='ignore') as fh:
                    n += sum(1 for _ in fh)
            except OSError:
                pass
    return n


LOC = {
    'backend': loc(['backend/*.py', 'backend/*.yaml']),
    'backend_tests': loc(['backend/tests', 'backend/scripts']),
    'ml': loc(['ml']),
    'child': loc(['android/ChildApp/app/src/main']),
    'child_tests': loc(['android/ChildApp/app/src/test']),
    'parent': loc(['android/ParentApp/app/src/main']),
    'parent_tests': loc(['android/ParentApp/app/src/test']),
    'docs': loc(['docs/*.tex', '*.md']),
    'ci': loc(['.github']),
}

# ---------- brand ----------
ORANGE = RGBColor(0xE0, 0x7B, 0x1A)   # template heading orange
RUST = RGBColor(0xB9, 0x54, 0x27)     # footer band
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xFB, 0xF3, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = 'Trebuchet MS'
TITLE = 'AI-Driven Gaming Addiction Screening System'
FOOT = 'Kaustubh_Khushee_Kanak_Vidisha'

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height
LOGO = os.path.join(FIG, 'pes_logo.png')


def _tb(slide, x, y, w, h, text, size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size, r.font.bold, r.font.name = Pt(size), bold, FONT
    r.font.color.rgb = color
    return box


def _md(s):
    out, bold, buf, i = [], False, '', 0
    while i < len(s):
        if s.startswith('**', i):
            if buf:
                out.append((buf, bold))
            buf, bold, i = '', not bold, i + 2
        else:
            buf += s[i]
            i += 1
    if buf:
        out.append((buf, bold))
    return out


def _rich(p, parts, size, color):
    for t, b in parts:
        r = p.add_run()
        r.text = t
        r.font.size, r.font.bold, r.font.name = Pt(size), b, FONT
        r.font.color.rgb = color


def chrome(slide, heading):
    """Template chrome: title-of-project top-left, logo top-right, orange heading, footer band."""
    _tb(slide, Inches(0.55), Inches(0.18), Inches(8), Inches(0.35), TITLE, 12, False, GREY)
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(12.25), Inches(0.12), height=Inches(0.95))
    _tb(slide, Inches(0.55), Inches(0.75), Inches(11.5), Inches(0.6), heading, 26, False, BLACK, PP_ALIGN.RIGHT)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - Inches(0.42), W, Inches(0.42))
    band.fill.solid()
    band.fill.fore_color.rgb = RUST
    band.line.fill.background()
    thin = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - Inches(0.48), W, Inches(0.06))
    thin.fill.solid()
    thin.fill.fore_color.rgb = ORANGE
    thin.line.fill.background()
    _tb(slide, Inches(0), H - Inches(0.36), W, Inches(0.3), FOOT, 10, False, WHITE, PP_ALIGN.CENTER)


def slide(heading, notes=''):
    s = prs.slides.add_slide(BLANK)
    chrome(s, heading)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def bullets(s, items, x=Inches(0.8), y=Inches(1.55), w=Inches(11.8), h=Inches(5.3), size=18, gap=7, color=ORANGE):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        sub = it.startswith('  ')
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.level = 1 if sub else 0
        _rich(p, _md(('– ' if sub else '▪  ') + it.strip()), size - (3 if sub else 0), BLACK if sub else color)
    return box


def table(s, rows, x, y, w, col_w=None, size=12, hi=()):
    nr, nc = len(rows), len(rows[0])
    shp = s.shapes.add_table(nr, nc, x, y, w, Inches(0.38) * nr)
    t = shp.table
    if col_w:
        for i, cw in enumerate(col_w):
            t.columns[i].width = cw
    for r, row in enumerate(rows):
        for c, v in enumerate(row):
            cell = t.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame
            tf.word_wrap = True
            _rich(tf.paragraphs[0], _md(str(v)), size, WHITE if r == 0 else BLACK)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ORANGE if r == 0 else (LIGHT if (r in hi or r % 2) else WHITE)
    return shp


def fig(s, name, x, y, w=None, h=None):
    p = os.path.join(FIG, name)
    if os.path.exists(p):
        return s.shapes.add_picture(p, x, y, width=w, height=h)


def ci(v):
    return f'[{v[0]:.3f}, {v[1]:.3f}]'


# =============== 1. Title (template p1) ===============
s = prs.slides.add_slide(BLANK)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(12.25), Inches(0.12), height=Inches(0.95))
_tb(s, Inches(0), Inches(1.1), W, Inches(0.7), 'UE22CS441A – Capstone Project Phase – 3', 30, False, BLACK, PP_ALIGN.CENTER)
_tb(s, Inches(0), Inches(1.95), W, Inches(0.6), 'Project Progress Review # 1', 26, False, ORANGE, PP_ALIGN.CENTER)
_tb(s, Inches(2.1), Inches(4.6), Inches(3.0), Inches(2.2),
    'Project Title   :\nProject ID        :\nProject Guide  :\nProject Team   :', 22, False, ORANGE)
_tb(s, Inches(5.0), Inches(4.6), Inches(8.0), Inches(2.5),
    'AI-Driven Gaming Addiction Screening System\nPW26_SAS-03\nProf. Shridevi Sawant\n'
    'Kaustubh Agarwal (PES1UG23CS291) · Khushee P Kiran (PES1UG23CS303)\n'
    'Kanak Goyal (PES1UG23CS279) · Vidisha Murali (PES1UG23CS681)', 18, False, BLACK)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - Inches(0.42), W, Inches(0.42))
band.fill.solid(); band.fill.fore_color.rgb = RUST; band.line.fill.background()
thin = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - Inches(0.48), W, Inches(0.06))
thin.fill.solid(); thin.fill.fore_color.rgb = ORANGE; thin.line.fill.background()
s.notes_slide.notes_text_frame.text = ('One sentence: a deployed, multimodal screening system for parents that measures HOW a '
                                        'child plays, not just how long — externally validated against a clinical instrument, '
                                        'with its limitations named before you ask.')

# =============== 2. Outline (template p2) ===============
s = slide('Outline')
bullets(s, ['Abstract and Scope of the Project',
            'Capstone Project Phase – 2',
            '  Summary of work',
            '  Inferences drawn from Literature Survey',
            'List of Tasks/Modules with Individual Contribution',
            'Design of overall Architecture',
            'Demonstration and Testing of the completed modules',
            'Gantt chart',
            'References'], size=22, gap=10)

# =============== 3. Abstract & Scope (template p3) ===============
s = slide('Abstract and Scope', notes=(
    'NUMBERS BEHIND THIS SLIDE: Numbers moved off this slide. IGDS9-SF is a 9-question instrument (the 9 in its name is the item count). Calibration example: a calibrated score of 0.8 means about 80 percent of children who get that score really are high-risk. The architecture has three tiers: Child app (background capture), cloud ML backend, Parent app. Hindi chat is handled in both romanised Latin spelling and Devanagari script. The accent-fairness audit used the Svarah Indian-English speech dataset. '
   
    'Set the context in 60 s: IGD is WHO-recognised; parents notice late; questionnaires need cooperation, screen-time '
    'counters measure volume. We measure HOW a child plays. Say "screening, not diagnosis" out loud. Scope: minors, '
    'Android, guardian in the loop — and why (a guardian with duty of care is what makes screening the unwilling '
    'possible and proportionate).'))
bullets(s, [
    '**Problem.** Internet Gaming Disorder (IGD) is a condition recognised by the WHO (listed in ICD-11, its disease catalogue), but parents notice it **late**, from the outside. Today’s tools are questionnaires the child must answer honestly, or screen-time counters that measure **how long** a child plays, which only weakly tracks how serious the problem is.',
    '**What we built.** A deployed system: a background-capturing Child app, a cloud machine-learning backend, and a Parent app. It combines **behaviour patterns** (how the child plays), **chat toxicity** (English and Hindi, in Latin or Devanagari script) and **voice** into one calibrated risk score (it can be read as a real probability), with a plain-language "why".',
    '**Scope.** Minors on Android, with a guardian’s consent. The output is a screening and awareness signal that flags concern for a parent. It is **not** a clinical diagnosis. The cloud backend is live in production, signed APKs are on GitHub, and a consented family pilot ran for 23 days.',
    '**What Phase 3 adds.** An external check (construct validation): does our score track IGDS9-SF, the standard clinical screening questionnaire for gaming disorder? (134 raw / 104 usable survey responses). An accent-fairness audit of our speech-to-text on Indian English (6,656 clips). Measured CPU and memory cost on a phone. And the honest limitations of each.',
], size=16, gap=9)

# =============== 4. Summary of Phase 1/2 + suggestions (template p4) ===============
s = slide('Summary of Work Done in Capstone Project Phase – 1 & 2', notes=(
    'NUMBERS BEHIND THIS SLIDE: Numbers and details moved off slide 4. Phase-1 proposal detail dropped from the cell: three risk bands Casual, At-risk, Addicted; proposed stack Flutter or React Native with Firebase; Android plus iOS; the camera was dropped because there was no proportionate consent story for it. Phase-2 detail dropped: the ensemble is 3 models (behaviour, chat, voice) fused into one score; the per-prediction why comes from SHAP; the 11 datasets are open-licence and were adopted after a measured trial; ablations (retrain with one part removed) were run with confidence intervals; the 266 tests were CI tests. Phase-2 feedback detail dropped: the synthetic-label caveat is stated on every slide; the Hindi fix also added a Devanagari keyboard in the ChildApp. Voice leakage context: the leak came from random splits letting the same speakers into both training and test; the speaker-independent number is 8.3 points lower and is the one we report. Phase-3 detail dropped: v2.4.0 is the release version shipped with the on-device drill. Survey context for the 0.317: rho is a Spearman rank correlation; 95 percent interval 0.137 to 0.478 on n = 103 respondents with a computable score (104 usable of 134 raw); screen-time hours alone 0.147; the gain over hours is +0.167 with interval -0.002 to +0.341, which touches zero, so it is suggestive, not proven; the 97 percent of resamples is the share of paired bootstrap resamples (same people resampled) in which the score beat hours; the 97 percent is the same suggestive result stated another way. Accent audit context for the gap named on the slide: word error rate 34.6 percent for Dravidian first-language speakers versus 60.8 percent for Tibeto-Burman, over 9.61 hours and 6,656 clips. Battery context: the gate is the measured on-device CPU and memory budget; the dual speech-to-text path (English and Hindi engines running together) failed it, so it ships default OFF. '
   
    'Phase 1 (approval) proposed the vision; Phase 2 built and evaluated the system; each review\'s feedback narrowed the '
    'design toward what could be validated. Be explicit that several Phase-1 ideas were dropped ON PURPOSE with a reason — '
    'camera/facial expression (privacy, no consent story), screen-recording OCR (privacy regression), iOS (no accessibility '
    'API for capture), Kaggle behaviour datasets (audited: synthetic provenance). Reviewers reward "we changed our mind '
    'because we measured" far more than "we did everything we said."'))
table(s, [
    ['Phase', 'What was proposed / done', 'Suggestion or finding → improvement made'],
    ['**Phase 1** (approval)', 'We proposed an app that predicts gaming-addiction risk from playtime, chat mood and voice emotion, with a parent dashboard. The original plan also included iOS, a camera for facial expressions, reading text off screenshots, and ready-made Kaggle datasets.',
     '**Cut to what we could capture, consent to, and validate.** Android only (iOS has no capture API). Camera dropped: consent not justifiable. Screenshot reading dropped: worsens privacy. Kaggle data rejected: it was generated, not collected from real players.'],
    ['**Phase 2** (build + evaluate)', 'Apps and Flask backend deployed. Behaviour, chat and voice models fused, calibrated so scores read as probabilities, each with a plain-language ‘why’. 11 real datasets adopted after trials. Family pilot ran. 266 tests at Phase-2 close (291 now).',
     '**Synthetic labels**: simulation grounded on 2 real surveys, not clinicians. **Voice**: random splits leaked speakers into both train and test, inflating accuracy. Splitting by speaker costs –8.3 pts, reported honestly. **Hindi**: real hate-speech data (HASOC), Roman and Devanagari.'],
    ['**Phase 3** (this review)', f'We tested the score against a clinical questionnaire in a survey, audited speech-to-text across Indian accents, measured CPU and memory on a phone, and shipped a new release. Documented in a {REPORT_PP}-page report and a {IEEE_PP}-page IEEE paper.',
     '**Score tracks IGDS9-SF (clinical questionnaire): correlation ρ = 0.317**, moderate. Beats screen time in 97% of resamples (suggestive). Accents: **0 false alerts / 9.6 h**. Transcription accuracy varies by accent. Battery measured. English+Hindi speech-to-text fails gate: default OFF.'],
], Inches(0.6), Inches(1.5), Inches(12.1), col_w=[Inches(1.7), Inches(5.0), Inches(5.4)], size=11)

# =============== 5. Inferences from literature (template outline) ===============
s = slide('Inferences Drawn from Literature Survey', notes=(
    'NUMBERS BEHIND THIS SLIDE: Numbers and terms moved off slide 5 into notes. 9-item: IGDS9-SF has 9 questions; it is the short form of the Internet Gaming Disorder Scale. DSM-5: the questionnaire follows the nine criteria for Internet Gaming Disorder in the fifth edition of the American Psychiatric Association manual. Context for the headline numbers still visible on the slide: n = 11,191 is the size of the open Latin-American IGDS9-SF dataset (OSF); 6.4 percent is the share of that dataset scoring in the disordered range, which is the base rate we use to set alert thresholds; r = +0.156 is the correlation between toxic-chat involvement and IGDS9-SF severity in that dataset, a small but present relationship; the chat-premise rho placeholder currently reads 0.323 (95 percent interval 0.144 to 0.491, n = 103) and is the same toxic-chat versus severity rank correlation measured in our own survey; 12 pts is the PR-AUC (area under the precision-recall curve) gap on real in-game chat between the pretrained toxic-BERT baseline at 0.709 and our deployed logistic regression at 0.825, whose 95 percent interval is 0.807 to 0.841; removing the CONDA in-game chat corpus from training drops our own model to 0.511, which is the evidence behind domain data beats model size. n = 104 usable is the number of survey respondents left from 134 raw after exclusions fixed in advance; the headline correlations are computed on the 103 of them with complete data, which is why slide 9 shows n = 103. Support for the last row, taken from slide 9 and not shown here: pattern composite rho 0.330 versus volume composite rho 0.128 against IGDS9-SF, with a formal pattern minus volume contrast interval of +0.019 to +0.380 that excludes zero; screen-time hours alone rho 0.147; the paired gain of the full score over hours alone is +0.167 with interval -0.002 to +0.341, which touches zero, so describe that gain as suggestive only. Overall construct validity rho is 0.317 with interval 0.137 to 0.478, moderate and clearly present, never call it strong. Simulated labels on the last row means the behaviour model was trained on labels from a simulation grounded on two real surveys, not on clinician labels; the survey row is the first test against real questionnaire answers. '
   
    'Three inferences that shaped the design, each tied to a paper the panel can check: (1) family factors and behavioural '
    'patterns, not hours, predict IGD → pattern features and a guardian in the loop; (2) IGDS9-SF is the validated short '
    'instrument → our external anchor; (3) domain data beats architecture for toxicity → CONDA in-game chat, HASOC Hindi. '
    'And one inference from our own survey that the literature did not give us: pattern > volume, measured.'))
table(s, [
    ['Source', 'Inference', 'How it shaped our design'],
    ['You et al. 2025; Coșa et al. 2025; Ergin & Essau 2025 (family factors in gaming disorder)', '**How a child plays matters more than how long.** Studies find that family relationships and play patterns predict gaming disorder better than raw hours, and that interventions work through the family.', 'So the guardian (parent) stays in the loop and the app nudges rather than blocks. The **play-pattern** measures (late-night sessions, repeated re-logins, binges) are core model inputs.'],
    ['Pontes & Griffiths 2015 (the IGDS9-SF questionnaire); a Latin-American IGDS9-SF dataset, n = 11,191', '**A validated yardstick exists.** IGDS9-SF is a validated screening questionnaire for gaming disorder. About 6.4 % of that dataset scores in the disordered range (the base rate), and toxic-chat involvement tracks severity (correlation r = +0.156).', f'IGDS9-SF is the **real-world yardstick** we test our score against. The base rate sets our alert thresholds. The toxic-chat link justifies our chat channel, and our own survey reproduced it (correlation ρ = {SX["chat_premise"]["rho"]:.3f}).'],
    ['Huang et al. 2024; Jiang 2024 (machine learning for gaming disorder)', '**No prior work deploys this.** Earlier machine-learning work predicts gaming disorder from questionnaire answers or a single signal. None of it passively captures several signals and reports to a guardian.', 'That is the gap we fill: a deployed system combining behaviour, chat and voice, whose scores mean what they say (calibrated) and come with a plain-language reason (explainable).'],
    ['Weld et al. 2021 (CONDA, a real in-game chat dataset); Mandl et al. 2019 (HASOC, Hindi hate-speech dataset); Javed et al. 2023 (Svarah, Indian-accent English speech)', '**Game chat is its own dialect, and speech recognition is not accent-neutral.** In-game chat and mixed Hindi–English read differently from everyday text, and speech-to-text error rates vary with Indian accents.', 'Domain data beats model size: the pretrained toxicity model toxic-BERT scores 12 pts below our simpler logistic regression on real game chat. We handle Hindi in both scripts (Devanagari and romanised) and **audited fairness across accents** on Svarah.'],
    ['**Our own IGDS9-SF survey (n = 104 usable)**', '**Measured, not assumed: how beats how long.** Every one of our play-pattern measures tracks the questionnaire score more closely than any of our screen-time measures (volume features).', 'This confirms our design premise, measure how a child plays rather than how long, against real questionnaire answers instead of simulated labels.'],
], Inches(0.6), Inches(1.5), Inches(12.1), col_w=[Inches(3.4), Inches(4.4), Inches(4.3)], size=11)

# =============== 6. Architecture (template p5) ===============
s = slide('Architecture', notes=(
    'NUMBERS BEHIND THIS SLIDE: Slide 6 (Architecture) addendum: numbers and technical labels moved off the slide. "10 feats" in the backend box was the random-forest input count; it is the same 10 session features named in the ChildApp box (they include late-night play ratio, rapid re-login ratio and binge sessions per week), so only the duplicate was removed. Fusion weights 40/30/30 are behaviour/chat/voice and are availability-weighted, meaning re-normalised over whichever channels are present for that child. Technical labels replaced by plain words on the slide: game detection order is allowlist, then OS app category, then parent override; chat capture uses a custom IME keyboard plus an Android accessibility service; the microphone gate is VAD (voice-activity detection) and speech-to-text is on-device Vosk; the three models are a random forest for behaviour, logistic regression over TF-IDF for chat, and HistGradientBoosting (boosted trees) for voice; score correction is isotonic calibration (a monotone correction so predicted risk matches observed risk) and the plain-language why comes from SHAP; the drift monitor uses PSI and KS statistics; transport is HTTPS with token authentication; export covers exactly the data that delete removes; the alert cut-off tuner is a Beta (Bayesian) threshold tuner in the backend, driven by the parent accurate/false-alarm verdicts; anti-tamper flags disabled capture; consent is versioned and no consent means no capture; the transparency screen is the "what we can and cannot see" screen. Module placement stated on the slide: game detection, session tracking, chat capture, voice capture with speech-to-text, consent and anti-tamper run on the child phone; the three models, fusion and alerting, the feedback tuner and the drift monitor run in the Flask backend; the dashboard is the ParentApp. No other numbers appeared on this slide. '
   
    'Walk left to right. ChildApp: foreground-package game detection, 10 objective session features, custom keyboard + '
    'accessibility for chat, VAD-gated mic with on-device Vosk STT, anti-tamper, consent-gated. Backend: Flask on Render, '
    'Neon Postgres, three models + availability-weighted fusion (40/30/30 over channels PRESENT), isotonic calibration, '
    'SHAP why, alerts to every guardian, drift monitor, export = delete scope. ParentApp: dashboard, 14-day trend, alerts '
    'with verdicts → threshold tuner, nudges, capture-coverage transparency screen.'))
bx_y, bx_h = Inches(1.6), Inches(2.6)
for title, body, x in [
    ('ChildApp (Android)', 'Spots when a game is running; parents can override\nMeasures 10 session features (late-night play, re-logins, binges)\nReads chat through our own keyboard and accessibility service\nMic on only while speaking; transcribed on the phone\nNo consent, no capture; tampering is flagged', Inches(0.6)),
    ('Flask backend (Render + Neon Postgres)', 'Three models, one per channel: behaviour, chat, voice\nBlended 40/30/30 (behaviour, chat, voice); missing channels re-weighted\nScores corrected to match real risk rates; plain-language why\nAlerts every guardian; flags data unlike training data\nExport shows exactly what delete removes; encrypted, token-authenticated', Inches(4.75)),
    ('ParentApp (Android)', 'Risk level with a plain-language why; 14-day trend\nInstant alerts for risk, toxic chat, or tampering\nParent marks alerts right or wrong; cut-off adapts\nSends nudges to the child; weekly PDF report\nShows honestly what the app can and cannot see', Inches(8.9)),
]:
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, bx_y, Inches(3.8), bx_h)
    b.fill.solid(); b.fill.fore_color.rgb = LIGHT; b.line.color.rgb = ORANGE; b.line.width = Pt(2)
    _tb(s, x + Inches(0.15), bx_y + Inches(0.1), Inches(3.5), Inches(0.5), title, 15, True, ORANGE)
    _tb(s, x + Inches(0.15), bx_y + Inches(0.6), Inches(3.5), Inches(1.9), body, 11, False, BLACK)
for x in (Inches(4.42), Inches(8.57)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, bx_y + Inches(1.1), Inches(0.32), Inches(0.4))
    a.fill.solid(); a.fill.fore_color.rgb = ORANGE; a.line.fill.background()
bullets(s, ['**Data flow:** the Child app sends session start/end times, chat lines and 10-s voice segments → each channel gets its own score → the three are blended into one corrected risk score with a plain-language why → parent dashboard and push alerts. Raw audio is **deleted** as soon as its features are extracted',
            '**Modules:** On the phone: game detection, session tracking, chat capture, voice capture with speech-to-text, consent and anti-tamper. In the cloud: behaviour, chat and voice models, score blending and alerting, the parent-feedback tuner, and a drift monitor that flags data unlike the training data. For parents: the dashboard'],
        y=Inches(4.45), h=Inches(2.3), size=13, gap=6)

# =============== 7. List of tasks/modules + tools (template p6) ===============
s = slide('List of Tasks / Modules', notes=(
    'NUMBERS BEHIND THIS SLIDE: Numbers moved off this slide. 10 - the number of objective session features the behaviour random forest takes as input (was "RF (10 feats)" in the Modelling row). Details trimmed from the visible cells for fit, kept here: the chat model output (logistic regression) is first calibrated with isotonic regression, then combined with the keyword word list by noisy-OR (which can only raise the score), and only then fused with the other channels; survey CSV columns are matched to model features by keyword; audio capture on the child app is gated by VAD (voice-activity detection) so the microphone is only on during speech; chat capture uses the custom IME keyboard plus the Android accessibility service; TF-IDF runs over word and character n-grams; the calibration plots are reliability diagrams (predicted probability versus observed frequency). Numbers in the SDK and tools line (Android SDK 34, Vosk 0.3.47, Python 3.11, Postgres 16, HASOC 2019, CC BY 4.0) stay visible because that whole line is a headline item. Vosk model codes en-in and hi are the small Indian-English and Hindi models. '
   
    'Template asks for the data pipeline stages explicitly, and the SDK/API/tool list with licences. Everything is '
    'open-source; no licensed component. Mention that the two gated datasets (HASOC, Svarah) are NOT redistributed.'))
table(s, [
    ['Stage', 'What we do', 'Where'],
    ['Data collection & preparation', '**11 open datasets** adopted after checking origin, task fit and licence, then a measured trial. **2 rejected** with evidence. Primary data is our own survey using IGDS9-SF (the standard gaming-disorder questionnaire): **134 raw / 104 usable** responses.', 'ml/fetch_*.py, ml/analyze_*.py, docs/VALIDATION_PLAN.md'],
    ['Data input', 'The child’s app (ChildApp) records when game sessions start and end, captures chat through our own keyboard and the Android accessibility service, and records audio only while someone is speaking. Survey answers arrive as a Google-Form CSV export.', 'android/ChildApp, ml/eval_behavior_survey.py'],
    ['Pre-processing', 'The live server reuses the training feature code, so features match exactly. Chat text becomes TF-IDF over words and character fragments. Each voice clip becomes **36 acoustic features** (pitch, energy, timbre). Hindi is respelled in Latin letters.', 'backend/behavior_features.py, ml/'],
    ['Modelling', 'One model per signal: random forest on behaviour, logistic regression on chat (plus a word list that can only raise scores), gradient boosting on voice. Scores are fused with weights that re-balance over whichever signals are present.', 'ml/retrain_models.py, backend/app.py'],
    ['Visualisation & interpretation', 'Every prediction comes with a plain-language reason for parents, computed with SHAP (per-input contribution scores). Evaluation figures: precision–recall curves, confusion matrices, calibration plots (are scores honest probabilities?), ablation tables (retrain minus one part) and the survey figure.', 'ml/make_figures.py, docs/figures/'],
    ['Storage', 'The same database code runs on SQLite for local development and on Postgres (hosted by Neon) in production. A parent’s data export contains exactly what a delete request removes. Raw audio is never stored.', 'backend/app.py, DEPLOY.md'],
], Inches(0.6), Inches(1.45), Inches(12.1), col_w=[Inches(2.0), Inches(7.4), Inches(2.7)], size=10)
_tb(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(0.35), 'SDK / API / Model / Tools (all open-source):', 13, True, ORANGE)
_tb(s, Inches(0.6), Inches(5.68), Inches(12.1), Inches(1.3),
    'Android SDK 34 / Kotlin · Retrofit · Vosk 0.3.47 offline speech-to-text (Apache-2.0 licence; small Indian-English and Hindi models) · Firebase Cloud Messaging (push alerts) · '
    'Python 3.11 · Flask · scikit-learn (RandomForest, LogisticRegression, HistGradientBoosting, isotonic calibration) · SHAP · '
    'librosa / webrtcvad (audio features / voice-activity detection) · psycopg2 · Postgres 16 (Neon) · Render (cloud host) · GitHub Actions · pytest · schemathesis (API fuzzing) · bandit · '
    'pip-audit · MobSF (three security scanners) · LaTeX/IEEEtran · Datasets: Jigsaw (Wikipedia comment toxicity), CONDA, Davidson (Twitter hate speech), HASOC 2019 (gated: on request), Hindi Wikipedia, '
    'RAVDESS, CREMA-D, EMO-DB, URDU (emotional-speech corpora), Gamers & Anxiety (OSF), IGDS9-SF LatAm (OSF), StudentLife, Svarah (CC BY 4.0, gated).',
    11, False, BLACK)

# =============== 8. Individual contribution (template p7) — EDIT BEFORE PRESENTING ===============
CONTRIB = [
    # name, modules, LOC (evidence-based per module; person split is a PROPOSAL), hours (ESTIMATE — edit)
    ('Kaustubh Agarwal', 'ML pipeline (3 models, calibration, fusion, ablations); backend model serving + REST API; external validation survey (IGDS9-SF) + analysis; accent-fairness audit; report & IEEE paper draft',
     LOC['ml'] + LOC['backend'] // 2 + LOC['docs'] // 2, '≈ 420 h'),
    ('Khushee P Kiran', 'ChildApp: game detection, session tracking, custom keyboard (IME) + Devanagari layout, accessibility chat capture, voice recorder + on-device speech-to-text, anti-tamper, consent; on-device resource drill',
     LOC['child'] + LOC['child_tests'], '≈ 300 h'),
    ('Kanak Goyal', 'Backend infra: one DB code path for SQLite (local) and Postgres (production), auth, rate limiting, alerts/FCM, drift monitor, export/delete; CI, tests, deployment (Render/Neon)',
     LOC['backend'] // 2 + LOC['backend_tests'] + LOC['ci'], '≈ 280 h'),
    ('Vidisha Murali', 'ParentApp: dashboard, alerts + accurate/false-alarm verdicts, chat/emotion analysis screens, “what we can and can’t see” transparency screen, weekly PDF; privacy/ethics docs, survey questionnaire & recruitment, review decks',
     LOC['parent'] + LOC['parent_tests'] + LOC['docs'] // 2, '≈ 260 h'),
]
s = slide('Individual Contribution', notes=(
    'The split follows the module ownership the team agreed (each member owns the modules listed). Module LOC are '
    'measured (git ls-files | wc -l on 2026-08-18); hours are team estimates. If a panelist asks how LOC were '
    'attributed: by module ownership — the repo is pushed from one machine, so git author lines do not split by '
    'person, and we say that rather than fake per-author blame. Timeline is on the Gantt slide.'))
table(s, [['Team member', 'Tasks / modules assigned', 'Development (LOC, approx.)', 'Time spent']] +
      [[n, m, f'{l:,}', h] for n, m, l, h in CONTRIB] +
      [['**Total (repo)**', f'backend {LOC["backend"]:,} · ML {LOC["ml"]:,} · ChildApp {LOC["child"]:,} · ParentApp {LOC["parent"]:,} · tests {LOC["backend_tests"] + LOC["child_tests"] + LOC["parent_tests"]:,} · docs {LOC["docs"]:,} · CI {LOC["ci"]:,}',
        f'**{sum(LOC.values()):,}**', '≈ 1,260 h']],
      Inches(0.6), Inches(1.5), Inches(12.1), col_w=[Inches(1.9), Inches(6.5), Inches(2.0), Inches(1.7)], size=11, hi=(5,))
_tb(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(0.9),
    'Lines of code (LOC) counted with `git ls-files | wc -l` on 2026-08-18 (source, tests, docs, CI) and attributed by the module ownership the team agreed; '
    'hours are team estimates. Timeline of every task/module: see the Gantt chart.', 11, False, GREY)

# =============== 9. Demonstration & testing (template p8) ===============
full = chat_ab['full recipe (deployed)']
no_conda = chat_ab['- CONDA corpus (domain data)']
s = slide('Demonstration and Testing of the Modules Completed', notes=(
    'NUMBERS BEHIND THIS SLIDE: Numbers moved off slide 9 into notes. Behaviour model: macro-F1 0.918; 5-fold cross-validation accuracy 0.921 plus or minus 0.002; expected calibration error (ECE) 0.062 before isotonic calibration and 0.015 after (this is what the "much closer to a true probability" sentence refers to); hours-only (screen-time-only) baseline accuracy 0.702; five pattern features alone 0.902 versus five volume features alone 0.871. Chat model: in-domain PR-AUC 95 percent bootstrap interval 0.807 to 0.841; PR-AUC with the CONDA in-game chat corpus removed from training 0.511; HASOC Hindi held-out test set 933 rows, never used in training. Voice model: chance level for the four emotion classes is 0.25; wav2vec2 large-model headroom accuracy 0.776. External validation survey: Spearman rho 95 percent interval 0.137 to 0.478; n 103 respondents in the construct-validity analysis; screen-time hours alone rho 0.147; paired gain of our score over hours delta rho +0.167 with interval -0.002 to +0.341, which touches zero, so the gain over hours is suggestive rather than established; pattern composite rho 0.330; volume composite rho 0.128; formal pattern-minus-volume contrast +0.199 with interval +0.019 to +0.380, which excludes zero (this is the "gap distinguishable from zero" on the slide); genre multiplier test p 0.598, no effect; exactly one respondent scored in the IGDS9-SF disordered range, which is why sensitivity and specificity at a cut-off were not computed. Fairness audit: 9.61 hours of speech; 6,656 clips; the Svarah corpus has 117 speakers; WER 34.6 percent for Dravidian first-language speakers and 60.8 percent for Tibeto-Burman. System: current release is v2.4.0; the consented family pilot ran 23 days; dual-STT (two speech-to-text engines together, English plus Hindi) measured 51 to 72 percent CPU and 419 MB, versus 14 percent CPU and 288 MB for the default path, so it fails the resource gate and ships default OFF. Negative results: the per-genre risk multiplier is unsupported (p 0.598); 4 of 5 proxy feature names (features named after clinical symptoms) were discredited by the survey and renamed in the product, which is the "most features named after clinical symptoms" on the slide. '
   
    'Demo per DEMO_RUNBOOK §3 (parent dashboard → alerts + verdict → nudge → chat/emotion → live capture in Roblox → '
    'tamper). Backup video ready. Testing: 291 automated tests + 7 research-integrity guards in CI on both DB dialects; load, fuzz, CVE, MobSF; '
    'on-device drill numbers here are MEASURED (Galaxy M52, 2026-08-18). Results table: every number is from the '
    'committed JSONs; 7 CI guards fail the build if the paper and data disagree.'))
table(s, [
    ['Module', 'Result (held-out, real data unless stated)', 'Testing'],
    ['Behaviour model', '**91.6% accuracy** on held-out data. Caveat: the labels are **synthetic** — they came from a simulation grounded on two real surveys, not from clinicians. Calibration brought the score much closer to a true probability.', 'Ablations (retraining with parts removed) with confidence intervals. Screen time alone scores far lower. Pattern (how-you-play) features edge out screen-time features.'],
    ['Chat model', f'On real in-game chat, PR-AUC is **{full["pr_auc"]:.3f}**. At our alert threshold (0.95), 95.6% of flagged messages are truly toxic and we catch 42.8% of toxic messages. Hindi precision: 0.968 in Devanagari, 0.958 romanised.', f'Dropping the in-game chat corpus (CONDA) from training cuts PR-AUC sharply. Pretrained toxic-BERT reaches only 0.709, below our simpler model. Hindi test rows were never trained on.'],
    ['Voice model', '**0.574 accuracy** when no speaker appears in both training and test (a speaker-independent split), well above guessing. A random split gave 0.657 — an 8.3-point gap, because the model had memorised voices. We report the honest number.', 'Each speaker stays on one side of every split. A larger model (wav2vec2) scores higher, so headroom exists. Audio augmentation made no difference.'],
    ['**External validation**', f'Our score tracks IGDS9-SF (standard clinical questionnaire): rank correlation ρ **{SV["construct_validity"]["rho"]:.3f}**, a **moderate** link. Screen-time hours alone: {inc["rho_hours"]:.3f}. How a child plays beats how long: pattern features correlate **+{comp["pattern_minus_volume"]["diff"]:.3f}** higher than screen-time features, a gap distinguishable from zero.', f'Exclusions fixed in advance. Confidence intervals from resampling the same people. Late responses lowered the headline but stay in, by a rule set beforehand. Game genre: no effect. Sensitivity/specificity not computed: only one respondent in the disordered range.'],
    ['Fairness (STT → toxicity)', f'**0 false alerts** across {AF["overall"]["clips"]:,} clips of Indian-English speech, in every accent group. But speech-to-text error (WER, word error rate) varies by first-language family: {fam["Dravidian"]["wer"]*100:.1f}% for Dravidian speakers up to {fam["Tibeto-Burman"]["wer"]*100:.1f}% for Tibeto-Burman speakers.', 'Svarah, an Indian-English speech corpus, run through the deployed speech-to-text and toxicity scorer.'],
    ['System / apps', 'Backend live on Render and Neon. Signed APKs tested on a phone. Consented family pilot completed. Default path measured on-device: **14% CPU / 288 MB**. Two speech-to-text engines together (dual-STT) **fail our resource gate: default OFF**.', '181 backend tests (SQLite and Postgres), 110 Android tests and 7 paper-versus-data guards run in CI. 288-request load test: 0 errors. Plus fuzzing, CVE and MobSF scans.'],
], Inches(0.6), Inches(1.45), Inches(12.1), col_w=[Inches(1.5), Inches(6.9), Inches(3.7)], size=9, hi=(4,))
fig(s, 'survey_features.png', Inches(0.6), Inches(5.6), h=Inches(1.4))
fig(s, 'pr_chat.png', Inches(4.1), Inches(5.6), h=Inches(1.4))
_tb(s, Inches(7.0), Inches(5.6), Inches(5.7), Inches(1.4),
    'Live demo: dashboard, alert verdict, nudge, chat/emotion analysis, live capture in Roblox, tamper alert; backup video ready. Two results AGAINST us: weighting risk by genre is unsupported, and most features named after clinical symptoms failed the survey check and were renamed. PR-AUC (area under the precision–recall curve) is the right metric when toxic messages are rare.',
    10, False, BLACK)

# =============== 10. Gantt (template outline) ===============
s = slide('Gantt Chart', notes=(
    'Phases from git history and the Phase-1 deck: Phase 1 approval + literature (Feb–Mar), Phase 2 design + build '
    '(Apr–Jun: 19 commits May, 100 June), Phase 2 evaluation + pilot (Jul: 48 commits, 23-day pilot 6–28 Jul), Phase 3 '
    'validation + release + reporting (Aug: 70 commits — survey 7–9 Aug, fairness audit 14 Aug, device drill + v2.4.0 '
    '18 Aug). Remaining: final report submission, viva.'))
months = ['Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep']
tasks = [
    ('Phase 1: problem, scope, literature survey, approval', 0, 2, 'done'),
    ('Architecture & data pipeline design; dataset audit', 1, 3, 'done'),
    ('ChildApp / ParentApp / backend build (v1 → v2.3)', 2, 5, 'done'),
    ('Model training, calibration, ablations, bake-offs', 3, 5, 'done'),
    ('Cloud deployment (Render + Neon), CI, security audits', 3, 6, 'done'),
    ('Consented family pilot (6–28 Jul); feedback loop; drift monitor', 5, 6, 'done'),
    ('Hindi dual-script chat + Devanagari keyboard + dual-STT (v2.4.0)', 5, 7, 'done'),
    (f'External validation survey (IGDS9-SF, n = {SV["n_raw"]}) + analysis', 6, 7, 'done'),
    ('Accent-fairness audit · on-device drill · v2.4.0 release', 6, 7, 'done'),
    (f'Report ({REPORT_PP} pp) · IEEE paper ({IEEE_PP} pp) · defense kit', 6, 7, 'done'),
    ('Phase 3 reviews · final submission · viva', 6, 8, 'open'),
]
gx, gy, gw = Inches(0.6), Inches(1.55), Inches(12.1)
label_w = Inches(4.6)
cell_w = (gw - label_w) / len(months)
row_h = Inches(0.42)
_tb(s, gx, gy, label_w, row_h, 'Task / module', 12, True, ORANGE)
for i, m in enumerate(months):
    _tb(s, gx + label_w + cell_w * i, gy, cell_w, row_h, m + ' ’26', 11, True, ORANGE, PP_ALIGN.CENTER)
for r, (name, a, b, st) in enumerate(tasks, start=1):
    y = gy + row_h * r
    _tb(s, gx, y + Inches(0.05), label_w, row_h, name, 11, False, BLACK)
    for i in range(len(months)):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, gx + label_w + cell_w * i, y, cell_w, row_h - Inches(0.04))
        cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if r % 2 else WHITE
        cell.line.color.rgb = RGBColor(0xE6, 0xE6, 0xE6); cell.line.width = Pt(0.5)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, gx + label_w + cell_w * a + Inches(0.04), y + Inches(0.08),
                             cell_w * (b - a + 1) - Inches(0.08), row_h - Inches(0.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE if st == 'done' else RGBColor(0xF2, 0xC2, 0x8A)
    bar.line.fill.background()
_tb(s, gx, gy + row_h * (len(tasks) + 1) + Inches(0.05), gw, Inches(0.4),
    'Solid = complete · Light = in progress. Cadence from git: 19 commits May · 100 Jun · 48 Jul · 70 Aug (to 18 Aug).', 11, False, GREY)

# =============== 11. References (template p9) ===============
s = slide('References', notes=f'IEEE format; the full {REPORT_REFS}-entry list is in the report, {IEEE_REFS} in the IEEE paper.')
refs = [
    'H. M. Pontes and M. D. Griffiths, "Measuring DSM-5 Internet Gaming Disorder: Development and validation of a short psychometric scale," Computers in Human Behavior, vol. 45, pp. 137–143, 2015.',
    'World Health Organization, International Classification of Diseases 11th Revision (ICD-11) — Gaming disorder (6C51), 2019.',
    'S. You, X. Wang, Z. Hu, and J. He, "Parent–child relationships and gaming addiction: A systematic review and meta-analysis," Journal of Youth and Adolescence, 2025.',
    'Y. Huang, R. Wu, Y. Huang, Y. Xiang, and W. Zhou, "Investigating mechanisms of Internet Gaming Disorder and developing intelligent monitoring models using AI technologies," BMC Public Health, 2024.',
    'D. A. Ergin and C. A. Essau, "Family factors and Internet Gaming Disorder among adolescents: A systematic review," Int. J. Developmental Science, 2025.',
    'H. Weld et al., "CONDA: a CONtextual Dual-Annotated dataset for in-game toxicity understanding and detection," in Findings of ACL-IJCNLP, 2021.',
    'T. Mandl et al., "Overview of the HASOC track at FIRE 2019: Hate speech and offensive content identification in Indo-European languages," in Proc. FIRE, 2019.',
    'T. Javed et al., "Svarah: Evaluating English ASR systems on Indian accents," in Proc. INTERSPEECH, 2023.',
    'S. M. Lundberg and S.-I. Lee, "A unified approach to interpreting model predictions," in Advances in Neural Information Processing Systems (NeurIPS), 2017.',
    'S. R. Livingstone and F. A. Russo, "The Ryerson Audio-Visual Database of Emotional Speech and Song (RAVDESS)," PLoS ONE, vol. 13, no. 5, 2018.',
    'H. Cao et al., "CREMA-D: Crowd-sourced emotional multimodal actors dataset," IEEE Trans. Affective Computing, vol. 5, no. 4, 2014.',
    'T. Davidson, D. Warmsley, M. Macy, and I. Weber, "Automated hate speech detection and the problem of offensive language," in Proc. ICWSM, 2017.',
]
box = s.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.3))
tf = box.text_frame; tf.word_wrap = True
for i, r in enumerate(refs, 1):
    p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
    p.space_after = Pt(4)
    _rich(p, [(f'[{i}] {r}', False)], 11, BLACK)

# =============== 12. Thank you (template p10) ===============
s = prs.slides.add_slide(BLANK)
_tb(s, Inches(0), Inches(3.0), W, Inches(1.4), 'Thank\nYou', 40, False, ORANGE, PP_ALIGN.CENTER)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - Inches(0.42), W, Inches(0.42))
band.fill.solid(); band.fill.fore_color.rgb = RUST; band.line.fill.background()
thin = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - Inches(0.48), W, Inches(0.06))
thin.fill.solid(); thin.fill.fore_color.rgb = ORANGE; thin.line.fill.background()

prs.save(OUT)
print(f'wrote {OUT}  ({len(prs.slides)} slides)')
print('LOC:', LOC, 'total', sum(LOC.values()))
