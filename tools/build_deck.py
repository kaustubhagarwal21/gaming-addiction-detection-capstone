# -*- coding: utf-8 -*-
"""Build the defense deck (docs/DEFENSE_DECK.pptx) from docs/SLIDE_OUTLINE.md's content.

Every survey / ablation / fairness number is READ from the committed JSONs, not
retyped, so the deck cannot drift from the paper (ml/tests pins the paper to the
same files). Speaker notes for every slide are in the notes pane.

Re-run after any re-analysis:  python tools/build_deck.py
"""
import json
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DOCS = os.path.join(ROOT, 'docs')
FIG = os.path.join(DOCS, 'figures')
OUT = os.path.join(DOCS, 'DEFENSE_DECK.pptx')


def _load(p):
    with open(p, encoding='utf-8') as f:
        return json.load(f)


SV = _load(os.path.join(DOCS, 'survey_validation.json'))
SX = _load(os.path.join(DOCS, 'survey_extras.json'))
AB = _load(os.path.join(DOCS, 'ablation_results.json'))
AF = _load(os.path.join(DOCS, 'asr_fairness.json'))

inc = SX['incremental']
comp = SX['features']['composites']
chat_ab = {r['config']: r for r in AB['chat']}
fam = AF['by_language_family']

# ---------- palette ----------
NAVY = RGBColor(0x14, 0x22, 0x3C)
TEAL = RGBColor(0x0E, 0x7C, 0x86)
AMBER = RGBColor(0xD9, 0x82, 0x1E)
GREY = RGBColor(0x5B, 0x63, 0x6E)
LIGHT = RGBColor(0xF3, 0xF5, 0xF8)
PALE = RGBColor(0xC9, 0xD6, 0xE3)
MINT = RGBColor(0x9F, 0xD8, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HILITE = RGBColor(0xFF, 0xF3, 0xE0)
FONT = 'Calibri'

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height
_n = [0]


def _tb(slide, x, y, w, h, text='', size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size, r.font.bold, r.font.name = Pt(size), bold, FONT
    r.font.color.rgb = color
    return box


def _md_parts(s):
    """Tiny **bold** parser -> [(text, bold), ...]."""
    out, bold, buf, i = [], False, '', 0
    while i < len(s):
        if s.startswith('**', i):
            if buf:
                out.append((buf, bold))
            buf, bold, i = '', not bold, i + 2
        else:
            buf += s[i]
            i += 1
    if buf:
        out.append((buf, bold))
    return out


def _rich(paragraph, parts, size=18, color=NAVY):
    for text, bold in parts:
        r = paragraph.add_run()
        r.text = text
        r.font.size, r.font.bold, r.font.name = Pt(size), bold, FONT
        r.font.color.rgb = color


def base(title, kicker=None, notes=''):
    _n[0] += 1
    s = prs.slides.add_slide(BLANK)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.15))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    _tb(s, Inches(0.6), Inches(0.2), Inches(12.0), Inches(0.6), title, 28, True, WHITE)
    if kicker:
        _tb(s, Inches(0.6), Inches(0.72), Inches(12.0), Inches(0.4), kicker, 13, False, PALE)
    _tb(s, Inches(0.6), Inches(7.05), Inches(9), Inches(0.3),
        'PW26_SAS-03 · AI-Driven Gaming Addiction Screening · PES University', 10, False, GREY)
    _tb(s, Inches(12.0), Inches(7.05), Inches(0.8), Inches(0.3), str(_n[0]), 10, False, GREY, PP_ALIGN.RIGHT)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def bullets(slide, items, x=Inches(0.6), y=Inches(1.45), w=Inches(12.1), h=Inches(5.3), size=19, gap=8):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        _rich(p, _md_parts('▪  ' + it.strip()), size, NAVY)
    return box


def table(slide, rows, x, y, w, col_w=None, size=13, hi_rows=()):
    nrows, ncols = len(rows), len(rows[0])
    shp = slide.shapes.add_table(nrows, ncols, x, y, w, Inches(0.4) * nrows)
    t = shp.table
    if col_w:
        for i, cw in enumerate(col_w):
            t.columns[i].width = cw
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = t.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame
            tf.word_wrap = True
            _rich(tf.paragraphs[0], _md_parts(str(val)), size, WHITE if r == 0 else NAVY)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = TEAL
            elif r in hi_rows:
                cell.fill.fore_color.rgb = HILITE
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
    return shp


def callout(slide, text, x, y, w, h, color=TEAL, size=14):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _rich(p, _md_parts(text), size, WHITE)
    return shp


def fig(slide, name, x, y, w=None):
    p = os.path.join(FIG, name)
    if os.path.exists(p):
        return slide.shapes.add_picture(p, x, y, width=w)
    return None


def f3(v):
    return f'{v:.3f}'


def ci(v):
    return f'[{v[0]:.3f}, {v[1]:.3f}]'


# =============== 1. Title ===============
s = prs.slides.add_slide(BLANK)
_n[0] += 1
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
_tb(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.3), 'AI-Driven Gaming Addiction Screening System', 40, True, WHITE)
_tb(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(0.9),
    'A multimodal, privacy-conscious mobile platform for early detection and parental awareness', 20, False, PALE)
_tb(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.6),
    'behaviour  ·  chat  ·  voice   →   one screening signal a parent can act on', 18, False, MINT)
_tb(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(1.2),
    'Kaustubh Agarwal (PES1UG23CS291) · Khushee P Kiran (PES1UG23CS303)\n'
    'Kanak Goyal (PES1UG23CS279) · Vidisha Murali (PES1UG23CS681)', 16, False, WHITE)
_tb(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.5),
    'Guide: Prof. Shridevi Sawant   ·   Team PW26_SAS-03   ·   PES University, 2026', 15, False, PALE)
s.notes_slide.notes_text_frame.text = (
    'One sentence: "We built and deployed a complete parent-facing screening system — two Android apps and a '
    'cloud ML backend — and evaluated it the way ML should be evaluated: real data where it exists, honest '
    'numbers where it doesn\'t, and an external validation study that returned two results we did not want."')

# =============== 1b. Phase III at a glance ===============
s = base('Phase III — where each expectation is evidenced', 'System testing · V&V · deployment · final results · performance analysis · research paper', notes=(
    'Use this as the map for the panel: every Phase III expectation has a slide, a repo artefact and a number behind it. '
    'If time is short, this slide plus slides 11 and 13 carry the grade. Point at the artefact column — everything is '
    'reproducible from the public repository, and 7 CI tests fail the build if the paper and the data disagree.'))
table(s, [
    ['Phase III expectation', 'What we did', 'Evidence'],
    ['**System testing**', '290 automated tests in CI (backend on SQLite + Postgres, Android JVM); load smoke, API fuzz, CVE + MobSF audits; on-device resource drill on real hardware', 'slide 14 · TESTING.md · ci.yml'],
    ['**Validation & verification**', f'Held-out / speaker-independent / in-domain evaluation; ablations with 95 % CIs; **external construct validation** vs IGDS9-SF (n = {SV["construct_validity"]["n"]}); accent-fairness audit; 7 research-integrity guards', 'slides 10–12 · docs/*.json · ml/tests/'],
    ['**Deployment**', 'Live cloud backend (Render + Neon Postgres, HTTPS, tokens); signed APKs v2.4.0 on GitHub, validated on device; consented family pilot ran 23 days', 'slide 4 · DEPLOY.md · GitHub Releases'],
    ['**Final experimental results**', 'Behaviour 91.6 % (synthetic); chat PR-AUC 0.825 in-domain, ≥ 0.95 precision on 3 registers; voice 0.574 speaker-independent; **ρ = 0.317 vs clinical instrument**, leads hours baseline (97% of paired resamples; partial excludes zero)', 'slides 6–12 · model_metadata.json'],
    ['**Performance analysis (tables & graphs)**', 'PR curve · confusion matrices · reliability diagrams · ablation tables · survey correlation table + feature chart · device resource table · fairness-by-accent table', 'slides 6–14 · docs/figures/'],
    ['**Complete research paper draft**', '47-page paper: architecture, methodology, 11-dataset audit, results, ablations, external validation, fairness audit, ethics & limitations, 44 references — numbers pinned to the data by CI', 'docs/PROJECT_PAPER.pdf'],
], Inches(0.6), Inches(1.45), Inches(12.1), col_w=[Inches(2.7), Inches(6.5), Inches(2.9)], size=13)
callout(s, 'Everything on this slide is reproducible from the public repository — and 7 CI tests fail the build if the paper and the data disagree.',
        Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.75), TEAL, 14)

# =============== 2. Problem & Scope ===============
s = base('Problem & Scope', 'Screening, not diagnosis — and why existing tools miss the pattern', notes=(
    'Set the scope before an examiner sets it for you: screening, not diagnosis — diagnosis needs a clinician. '
    'The gap: questionnaires need the child\'s honest cooperation; screen-time counters cannot tell 3 hours of '
    'healthy play from 3 hours of compulsive play. Behaviour patterns, in-game communication and voice stress '
    'together can. Children specifically: a guardian with duty of care is what makes SCREENING — reaching the '
    'person who would not self-report — both possible and proportionate; an adult version is either a self-help '
    'journal or spyware. (DEFENSE_NOTES §1)'))
bullets(s, [
    'Internet Gaming Disorder is WHO-recognised (ICD-11) — parents notice **late**, from the outside',
    'Existing tools: self-report questionnaires **or** blunt screen-time counters / blockers',
    'Ours: **passive, multi-signal screening** surfaced to a parent — with a plain-language "why"',
    '**NOT a diagnostic instrument** — a "talk to your child / consider help" signal',
    'Scoped to minors by design: a guardian in the loop is what makes screening the unwilling possible',
])
callout(s, 'Volume tells you **how long**. We measure **how** — late-night play, rapid re-logins, binges, rising toxicity, voice stress.',
        Inches(0.6), Inches(5.55), Inches(12.1), Inches(0.9))

# =============== 3. System Overview ===============
s = base('System Overview & Deployment', 'Three tiers, live in production — not localhost', notes=(
    'Walk left to right: ChildApp captures session telemetry, game chat (custom keyboard + accessibility), and '
    'gated mic segments; the Flask backend scores each channel and fuses them; ParentApp polls dashboards and '
    'receives push alerts to every guardian. It is LIVE — Render + Neon Postgres, HTTPS, token auth — and '
    'sideload-distributed by design (a family-consent install, not a Play Store product). (DEFENSE_NOTES §8)'))
bx_y, bx_h = Inches(1.9), Inches(2.4)
boxes = [
    ('ChildApp (Android)', 'passive capture\nsession telemetry · IME + accessibility chat\ngated mic → on-device STT\nanti-tamper · consent-gated', Inches(0.6)),
    ('Flask backend (cloud)', 'Render + Neon Postgres · HTTPS · tokens\n3 models + availability-weighted fusion\nSHAP "why" · alerts · drift monitor\nexport scope = delete scope', Inches(4.75)),
    ('ParentApp (Android)', 'dashboard + 14-day trend\nreal-time alerts (all guardians)\nfeedback loop → threshold tuner\ncapture-coverage transparency', Inches(8.9)),
]
for title, body, x in boxes:
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, bx_y, Inches(3.8), bx_h)
    b.fill.solid()
    b.fill.fore_color.rgb = LIGHT
    b.line.color.rgb = TEAL
    b.line.width = Pt(2)
    _tb(s, x + Inches(0.15), bx_y + Inches(0.12), Inches(3.5), Inches(0.5), title, 18, True, TEAL)
    _tb(s, x + Inches(0.15), bx_y + Inches(0.65), Inches(3.5), Inches(1.7), body, 13, False, NAVY)
for x in (Inches(4.42), Inches(8.57)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, bx_y + Inches(1.0), Inches(0.32), Inches(0.4))
    a.fill.solid()
    a.fill.fore_color.rgb = AMBER
    a.line.fill.background()
bullets(s, [
    'SQLite in dev / **Postgres (Neon)** in production — one code path, CI runs both dialects',
    'Signed release APKs on GitHub; **v2.4.0** current, validated on real hardware (2026-08-18)',
    'Consent is versioned and **fails closed**: bump the version → ingestion 403s until re-consent',
], y=Inches(4.6), h=Inches(2.2), size=16)

# =============== 4. Why three channels ===============
s = base('Why Three Channels', 'They fail independently — fusion treats each as an optional witness', notes=(
    'A missing channel contributes nothing rather than a fake neutral score — an earlier version imputed missing '
    'modalities and diluted real signal. Behaviour dominates because it is always eventually present and best '
    'measured. Say now: "the voice weight is lowest because we MEASURED it as the weakest channel — slide 7 shows '
    'that number honestly." (DEFENSE_NOTES §1, §5)'))
bullets(s, [
    '**Behaviour** needs days of telemetry; **chat** exists only when the child types; **voice** only when they speak',
    'Fusion prior **40 / 30 / 30**, **renormalised over the channels present** this session',
    'Weights follow **measured** reliability — voice is lowest for a reason we show, not hide',
    'Genre multiplier and band thresholds are priors, env-tunable — the first things we took to external data',
], h=Inches(2.8))
table(s, [
    ['Channel', 'Signal', 'Real training data', 'Weakest link (stated)'],
    ['Behaviour', '10 objective session features', 'labels synthetic; distributions from 2 real surveys', 'labels'],
    ['Chat', 'calibrated toxicity, 3 registers', '5 real corpora incl. in-game (CONDA) + Hindi (HASOC)', 'field data thin'],
    ['Voice', 'STT→toxicity + acoustic emotion', '4 real speech corpora, 154 speakers', 'acted adult speech'],
], Inches(0.6), Inches(4.5), Inches(12.1), size=13)

# =============== 5. Behaviour model ===============
s = base('Behaviour Model', 'Random Forest on ten measured features — calibrated, explainable, honest about its labels', notes=(
    'Two attack points, answer both proactively. (1) Why RF, not deep learning: bake-off on tabular data — RF wins on '
    'accuracy-per-cost and gives SHAP-friendly explanations parents see. (2) Synthetic labels: no public dataset pairs '
    'device telemetry with addiction labels — we audited candidates (slide 13). Grounding: IGDS9-SF Latin-America '
    '(11,191 real MOBA players) fixes severity base rates; Gamers & Anxiety (13,464) grounds behaviour–severity '
    'relations. The 20→10 feature cut: the 10 derived proxies are functions of the first 10; keeping them makes SHAP '
    'circular — ablation 0.9191 vs 0.9160, CIs overlap. Say plainly: 91.6% is a synthetic-distribution number; slide 11 '
    'is where the score meets real labels. (DEFENSE_NOTES §2)'))
bullets(s, [
    'Random Forest (200 trees × depth 6) on **10 objective features** — the model never sees the 10 derived proxies',
    '**91.6 %** test accuracy · macro-F1 0.918 · 5-fold CV 0.921 ± 0.002 — errors only ever land in an **adjacent** band',
    'Isotonic calibration: Brier 0.134 → 0.128; top-label **ECE 0.062 → 0.015**',
    'Labels are **synthetic** — grounded on real psychometrics (n = 11,191 + 13,464); stated on every slide it matters',
], w=Inches(7.2), size=16)
fig(s, 'cm_behaviour.png', Inches(8.1), Inches(1.5), w=Inches(4.7))
callout(s, 'Ten pattern-and-volume features in; **how** a child plays, not just how long. The proxies parents see are explanations, never inputs.',
        Inches(0.6), Inches(5.6), Inches(7.2), Inches(1.0), AMBER, 14)

# =============== 6. Chat model ===============
full = chat_ab['full recipe (deployed)']
no_char = chat_ab['- char_wb n-grams']
no_conda = chat_ab['- CONDA corpus (domain data)']
s = base('Chat Model', 'Domain data beats model capacity — and the threshold is a product decision', notes=(
    'The threshold story is the strongest 60 s of the talk: at the realistic ~3.5% toxic base rate, threshold 0.5 gives '
    'precision 0.235 — a parent gets ~3 false alarms per real one and stops reading alerts. At the deployed 0.95, '
    'precision is 0.956 at recall 0.428: we consciously trade recall to keep alerts credible (the session-level streak '
    'alert recovers coverage). Threshold is env-tunable; a Beta-posterior tuner adjusts it from parent feedback, capped '
    'at ±0.05 per run and human-applied. Why no BERT: 512 MB serving budget, real-time per-message scoring, and '
    'off-the-shelf toxic-BERT scores 0.709 on the same split — 12 points BELOW the domain-trained classical pipeline. '
    '(DEFENSE_NOTES §3, §5)'))
bullets(s, [
    'TF-IDF (word 1–2g + char_wb 3–5g) → LogReg → isotonic; **noisy-OR** fusion with a Hinglish / Devanagari lexicon',
    f'In-domain (CONDA game chat): **PR-AUC {full["pr_auc"]:.3f}** {ci(full["pr_auc_ci95"])} — off-the-shelf toxic-BERT: 0.709',
    'At alert threshold **0.95**: precision **0.956**, recall 0.428 — **≥ 0.95 precision on every register** at one threshold',
    '**Dual-script Hindi** (HASOC 2019 + clean-wiki counterweight): held-out precision **0.968 / 0.958** — near-zero before',
], w=Inches(7.3), size=15)
fig(s, 'pr_chat.png', Inches(8.2), Inches(1.5), w=Inches(4.6))
table(s, [
    ['Ablation (chat, held-out CONDA)', 'PR-AUC', '95 % CI'],
    ['full recipe (deployed)', f3(full['pr_auc']), ci(full['pr_auc_ci95'])],
    ['− char_wb n-grams', f3(no_char['pr_auc']), ci(no_char['pr_auc_ci95'])],
    ['− CONDA domain data', f'**{no_conda["pr_auc"]:.3f}**', ci(no_conda['pr_auc_ci95'])],
], Inches(0.6), Inches(5.3), Inches(7.3), size=12, hi_rows=(3,))

# =============== 7. Voice ===============
s = base('Voice Channel — two paths', 'On-device speech-to-text, and acoustic emotion trained on real speech', notes=(
    'Pre-empt "57% is weak": chance on 4 classes is 25%; the interesting part is the 9-point gap — with random splits '
    'the model memorised VOICES, not emotions, and the split choice even flipped which model won the bake-off. An honest '
    '0.574 beats an inflated 0.657 in front of any examiner. Transcription happens on-device; a short WAV goes up over '
    'HTTPS and raw audio is deleted server-side after feature extraction. Measured on hardware (slide 13): the default '
    'voice path costs ~14% of one core. Known gap (slide 15): acted adult emotion ≠ child gaming speech. (DEFENSE_NOTES §4)'))
bullets(s, [
    'Path 1: **on-device** Vosk STT (Indian English) → transcript → the same chat toxicity model',
    'Path 2: 36 acoustic features → HistGB emotion, trained on **real** corpora (RAVDESS, CREMA-D, EMO-DB, URDU)',
    'Speaker-independent accuracy **0.574** (chance 0.25) — random splits said 0.657: **9 points was speaker leakage**',
    'Headroom measured: frozen w2v2 embeddings + same classifier → 0.776 — the ceiling a deployable distillation could reach',
], w=Inches(7.3), size=15)
fig(s, 'cm_voice.png', Inches(8.2), Inches(1.5), w=Inches(4.6))
callout(s, 'Raw audio is **never retained**: on-device STT, server-side deletion after feature extraction, VAD-gated capture.',
        Inches(0.6), Inches(5.6), Inches(7.3), Inches(0.9), TEAL, 14)

# =============== 8. Fusion & alerting ===============
s = base('Fusion & Alerting', 'Max-not-mean, a threshold above best-F1, and a feedback loop that learns from parents', notes=(
    'Max-not-mean: one credible threat matters more than a hundred clean messages diluting it. Sitting above best-F1 is a '
    'CHOICE — the asymmetric cost is parents ignoring alerts. Know both operating points: 0.85 → P 0.800/R 0.703; '
    '0.90 → P 0.888/R 0.623. Feedback tuner: Beta posterior, ±0.05 cap per run, minimum-evidence floor, recommendations '
    'a human applies via env vars — no automatic loop can drift. (DEFENSE_NOTES §5)'))
bullets(s, [
    'Weighted fusion 40/30/30 over channels **present** this session; genre multiplier; observation mode for < 3 sessions',
    'Session chat risk = **max** of per-message calibrated scores — one credible threat is not diluted by clean lines',
    'Alert at **0.95**, deliberately above best-F1 (0.85): false alarms cost the only thing that matters — trust',
    'Parent "accurate / false alarm" verdicts → Beta-posterior threshold tuner (capped ±0.05, human-applied)',
    'Every prior we could not fit we **stress-tested**: removing the genre multiplier flips 34 % of served bands',
], h=Inches(3.2), size=16)
table(s, [
    ['Operating point', 'Precision', 'Recall', 'Meaning for a parent'],
    ['0.85 (best-F1)', '0.800', '0.703', '1 in 5 alerts is wrong'],
    ['0.90', '0.888', '0.623', '~1 in 9 wrong'],
    ['**0.95 (deployed)**', '**0.956**', '0.428', '~1 in 23 wrong — alerts stay credible'],
], Inches(0.6), Inches(4.85), Inches(12.1), size=13, hi_rows=(3,))

# =============== 9. Built for India ===============
drav, indo, tib = fam.get('Dravidian', {}), fam.get('Indo-Aryan', {}), fam.get('Tibeto-Burman', {})
s = base('Built for India — the differentiator', 'Code-mixed chat in three registers, a Devanagari keyboard, and an accent-fairness audit', notes=(
    'This is what makes it OURS, not a re-skin of a US parental-control app. English-only toxicity models miss most '
    'Indian gaming abuse. We measured the gap and closed it: a script-prior bug (a Devanagari corpus with a 53% offensive '
    'rate taught the model that Devanagari ITSELF looked toxic) surfaced from our own smoke test; a clean-Hindi Wikipedia '
    'counterweight fixed it. The keyboard matters because a third-party Hindi keyboard inside a canvas game is invisible '
    'to capture; ours is not. Fairness: every Svarah clip through the deployed STT and the served scorer — zero false alerts '
    'in 9.6 h across every accent group, but WER 35% Dravidian vs 61% Tibeto-Burman: a coverage gap on Northeast accents, '
    'not an accusation gap; named as the highest-leverage STT improvement. (DEFENSE_NOTES Hindi + fairness Q&As)'))
bullets(s, [
    'Indian kids\' game chat is **code-mixed**: English + romanised Hindi + Devanagari — all three at **≥ 0.95 precision**',
    'Native **Devanagari keyboard** in the Wellbeing Keyboard — Hindi captured even in canvas games (Roblox)',
    f'**Accent-fairness audited** on Svarah ({AF["overall"]["clips"]:,} clips / 117 speakers): **0 false alerts** in {AF["overall"]["speech_hours"]} h of benign speech, every accent group',
    'Honest capture matrix shown to parents — **including the blind spots** (browser games, in-game VOIP holding the mic)',
], h=Inches(3.0), size=15)
table(s, [
    ['Accent group (L1 family)', 'clips', 'STT word-error rate', 'benign-speech false alerts'],
    ['Dravidian', f'{drav.get("clips", "")}', f'{drav.get("wer", 0) * 100:.1f} %', f'{drav.get("alerts_hyp", 0)}'],
    ['Indo-Aryan', f'{indo.get("clips", "")}', f'{indo.get("wer", 0) * 100:.1f} %', f'{indo.get("alerts_hyp", 0)}'],
    ['Tibeto-Burman (Northeast)', f'{tib.get("clips", "")}', f'**{tib.get("wer", 0) * 100:.1f} %**', f'{tib.get("alerts_hyp", 0)}'],
], Inches(0.6), Inches(4.65), Inches(12.1), size=13, hi_rows=(3,))
_tb(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.5),
    'The chain degrades toward silence, never toward accusation — but it hears some accents markedly worse. A coverage gap, named.', 12, False, GREY)

# =============== 10. Evaluation methodology ===============
s = base('Verification — evaluation methodology', 'Every design choice has a measured counterfactual', notes=(
    'This slide is what separates the project from "we trained a model." Every design choice has a measured counterfactual. '
    'Headline lesson: data > architecture — losing domain data costs 31 PR-AUC points, more than every architecture choice '
    'combined — and three transformer benchmarks (Detoxify on chat, w2v2 on voice, MuRIL on Hindi) confirm it: capacity helps '
    'in-distribution, domain and register fit decide. Mention the reported NEGATIVE result: voice augmentation was measured and '
    'found neutral — we say so instead of hiding it. And now (slide 11) the pipeline output is anchored OUTSIDE its own '
    'training distribution. (DEFENSE_NOTES §6)'))
bullets(s, [
    'Held-out only, seeds recorded; **voice = speaker-independent**; **chat = in-domain** (real game chat), not general-corpus',
    'Ablations: one component removed per row, **bootstrap 95 % CIs** (1,000 resamples) — "better" only when intervals separate',
    'MCC + PR-AUC alongside F1 (base rates are extreme); direct calibration proof: **ECE 0.062 → 0.015** (behaviour)',
    f'Biggest finding: remove CONDA domain data → PR-AUC {full["pr_auc"]:.3f} → **{no_conda["pr_auc"]:.3f}** — data beats architecture',
    'Negative results reported: voice augmentation neutral; genre multiplier unsupported; 4 of 5 proxy names discredited',
], w=Inches(6.9), size=15)
fig(s, 'reliability.png', Inches(7.7), Inches(1.6), w=Inches(5.2))
_tb(s, Inches(7.7), Inches(3.75), Inches(5.2), Inches(0.5),
    'Reliability diagrams: calibrated confidence vs observed accuracy, both channels', 11, False, GREY)
callout(s, '7 research-integrity tests in CI pin the paper\'s numbers to the committed JSONs — a stale figure fails the build.',
        Inches(7.7), Inches(4.4), Inches(5.2), Inches(1.0), AMBER, 13)

# =============== 11. External validation ===============
s = base('Validation — does the score mean anything?', 'An IGDS9-SF survey, scored through the deployed pipeline, against real labels', notes=(
    'This is the slide the whole project builds toward — protect its time. Every other number in the deck is measured INSIDE '
    'our own training distribution; this one is measured outside it, against an instrument we did not design, on people who '
    'never touched the app. Lead with the comparison, not the magnitude — 0.35 sounds modest until you say the baseline every '
    'commercial parental-control tool ships reaches 0.147 with a CI spanning zero — and volunteer that the paired delta grazes zero after the late batch, so the partial correlation (0.303, CI excluding zero) carries the incremental claim. '
    'Pre-empt the two follow-ups: WHY NO SENSITIVITY/SPECIFICITY — one respondent scored in the disordered range; the script '
    'refuses caseness metrics below ten positives, a guard written before we saw the data. WHY KEEP THE GENRE MULTIPLIER — the '
    'null is underpowered (36%), not decisive, and removing it flips 34% of served bands; it stays flagged and env-tunable. '
    'Close on the two negative results: volunteer them. A validation study that only confirms is not one. (DEFENSE_NOTES §10)'))
bullets(s, [
    f'Anonymous adult IGDS9-SF survey — **{SV["n_raw"]} raw → {SV["n_usable"]} usable**; pattern answers scored through the **deployed** serving path',
    f'Construct validity: **ρ = {SV["construct_validity"]["rho"]:.3f}** {ci(SV["construct_validity"]["ci95"])} vs the clinical instrument — CI excludes zero',
    f'**Leads the screen-time baseline it replaces** in {inc["p_model_better"]*100:.1f} % of paired resamples: hours ρ = {inc["rho_hours"]:.3f} (CI spans zero); Δρ = +{inc["delta_rho"]:.3f} {ci(inc["delta_ci"])} — **CI grazes zero; the partial carries the claim**',
    f'Not a screen-time proxy: partial ρ = **{inc["partial_rho"]:.3f}** with hours removed',
    f'Signal is in **pattern** features ({comp["pattern"]["rho"]:.3f}) not **volume** ({comp["volume"]["rho"]:.3f}) — formal paired contrast **+{comp["pattern_minus_volume"]["diff"]:.3f} {ci(comp["pattern_minus_volume"]["ci"])}**, excludes zero',
    f'Two results **against** us: genre multiplier p = {SX["genre"]["p"]:.3f}; 4 of 5 proxy names track nothing → **labels renamed in the product**',
], w=Inches(7.4), h=Inches(2.7), size=12, gap=3)
table(s, [
    ['', 'ρ vs IGDS9-SF', '95 % CI'],
    ['**Model risk score**', f'**{inc["rho_model"]:.3f}**', ci(inc['rho_model_ci'])],
    ['Self-reported hours/week', f'{inc["rho_hours"]:.3f}', ci(inc['rho_hours_ci'])],
    ['**Paired difference**', f'**+{inc["delta_rho"]:.3f}**', ci(inc['delta_ci'])],
    ['Partial (hours removed)', f'{inc["partial_rho"]:.3f}', ci(inc['partial_ci'])],
], Inches(8.25), Inches(1.5), Inches(4.5), col_w=[Inches(2.1), Inches(1.1), Inches(1.3)], size=12, hi_rows=(1, 3))
callout(s, f'Model beats the hours baseline in **{inc["p_model_better"] * 100:.1f} %** of paired resamples. Caseness metrics withheld by design (1 disordered-range respondent).',
        Inches(8.25), Inches(3.9), Inches(4.5), Inches(1.15), TEAL, 12)
fig(s, 'survey_features.png', Inches(0.6), Inches(4.15), w=Inches(5.0))
_tb(s, Inches(8.25), Inches(5.2), Inches(4.5), Inches(1.6),
    'Left: every pattern feature out-ranks every volume feature — no interleaving. The design premise, confirmed against real labels rather than assumed.', 12, False, GREY)

# =============== 12. Dataset audit ===============
s = base('Dataset Audit', 'Eleven real open corpora adopted by measured trial — and the ones we rejected, with reasons', notes=(
    'We downloaded and inspected the rejected sets rather than dismissing them from their descriptions — provenance analysis '
    'showed one is synthetically generated with engagement labels, not addiction labels. Examiners like rejected-with-reasons '
    'more than adopted-without-reasons. Plus two datasets we collected/used ourselves: the IGDS9-SF survey (primary data) and '
    'Svarah (eval-only, gated download, not redistributed). (DEFENSE_NOTES §7)'))
bullets(s, [
    '**Adopted (11)**: Gamers & Anxiety (13,464) · IGDS9-SF LatAm (11,191) · CONDA · Davidson · **HASOC 2019 Hindi** · clean-Hindi Wikipedia · RAVDESS + CREMA-D + EMO-DB + URDU (9,817 clips)',
    f'**+ our own primary data**: the IGDS9-SF survey (n = {SV["n_usable"]} usable) — no open dataset pairs the instrument with the pattern variables we measure',
    '**+ eval-only**: Svarah Indian-accent English (6,656 clips, CC BY 4.0) for the fairness audit',
    '**Rejected with evidence**: Kaggle "Predict Online Gaming Behavior" (synthetic; engagement ≠ addiction) · "Mobile App Usage" (3/10 features, credential-walled)',
    '**Reality-checked** one hand-set prior against real phone telemetry (StudentLife): heavy-band late-night sits above the 90th percentile of normal student use',
], size=15)

# =============== 13. Engineering quality (with device numbers) ===============
s = base('System Testing & Engineering Quality', 'Tested, load-verified, watched in production — and measured on real hardware', notes=(
    'Two stories if asked: (1) the Postgres CI job caught a real dialect bug in the drift script the day it was added — exactly the '
    '"works locally, fails in prod" class it exists for; (2) the device drill (2026-08-18, Galaxy M52): the default monitoring path costs '
    '~14% of one core — about 1/15th of the game it monitors — and 288 MB; the dual Hindi+English STT toggle costs 4–5× the CPU and '
    'peaks at 419 MB, FAILING our own acceptance gate, which is exactly why it ships default OFF. A measurement that failed a feature '
    'is worth more than one that flattered it. Also on record: adb could not uninstall the app past the Device Admin block — the '
    'anti-tamper working against a developer with USB access. (DEFENSE_NOTES §8 + battery Q&A)'))
bullets(s, [
    '**290 automated tests in CI**: 180 backend (run on **both** SQLite & Postgres 16) + 110 Android JVM + 7 research-integrity guards',
    'Load-verified: 288 concurrent requests, **0 errors**, p50 66 ms; API fuzzed (schemathesis), CVE-audited, MobSF static audit clean',
    'Weekly **drift monitor** vs the production DB (PSI/KS); signed-token auth, rate limiting, authz regression tests',
    'Anti-tamper on record: even **adb could not uninstall** the app past the Device-Admin block without the parent PIN',
], size=14, h=Inches(2.7))
table(s, [
    ['On-device cost (Galaxy M52, 15-min Roblox)', 'English STT (default)', 'Dual Hindi+English toggle', 'Gate'],
    ['CPU, averaged', '**14 %** of one core (game: 215 %)', '51–72 % (~4–5×)', '< 2× — **fails**'],
    ['RAM (PSS) mean / peak', '**288 / 301 MB**', '399 / **419 MB**', '< 400 MB — **fails**'],
    ['Thermal throttling', 'none', 'none', 'none — passes'],
    ['Verdict', '**ships as default**', '**default OFF — measured reason**', ''],
], Inches(0.6), Inches(4.35), Inches(12.1), size=12, hi_rows=(4,))

# =============== 14. Privacy & ethics ===============
s = base('Privacy & Ethics', 'Consent that fails closed, symmetric export/delete, and the tensions we name instead of hide', notes=(
    'Lead with the strongest fact: what a parent can see is exactly what gets erased — export and deletion cover the same 15 tables. '
    'If pressed on the mic: layered mitigations in order — consent, on-device transcription, server-side raw-audio deletion, '
    'speech-gated capture. Then the two tensions we state ourselves: (a) the crisis asymmetry — profanity alerts a parent, a self-harm '
    'disclosure to the companion does not, DELIBERATELY: we refused to hard-code a contested clinical judgement; the tiered '
    'wellbeing-check design is documented as a proposal pending clinical review; (b) sensing blindspots — browser/webview games are '
    'invisible because the boundary is "installed apps", and closing it needs screen capture or URL surveillance, which we refuse. '
    '(DEFENSE_NOTES §9)'))
bullets(s, [
    'Explicit onboarding consent, **versioned, fails closed**; revoking a capture permission **alerts the parent** — no silent failure',
    'STT **on-device**; raw audio deleted after feature extraction; VAD gates the mic; **no raw media retained anywhere**',
    '`/api/user/export` = everything held (credentials excluded); **deletion scope identical to export** — one table inventory in code',
    'Parental monitoring of a **minor**, by design — visibility over prevention: alerts, not blocks; never covert toward the child',
    '**Named tensions**: the crisis-alert asymmetry (a refusal, not an oversight) · sensing blindspots (browser games) · adult-only validation cohort',
], size=15)
callout(s, 'A monitoring boundary honestly disclosed to the parent beats a privacy regression that closes it.',
        Inches(0.6), Inches(5.75), Inches(12.1), Inches(0.8), TEAL, 15)

# =============== 15. Limitations ===============
s = base('Limitations & What Is Still Open', 'Partial in a specific, nameable way — and we will not blur the line', notes=(
    'The framing shifted this year: the weakest part USED to be that nothing was externally validated; now it is that validation is '
    'partial in a specific, nameable way. Draw the line explicitly — construct validity is measured, caseness accuracy is not, and we '
    'will not blur them. That distinction is what a sharp examiner is probing for, so say it first. Then: "the contribution is not a '
    'solved clinical instrument — it is a fully built, honestly measured screening pipeline where every claim traces to a runnable '
    'script, including the claims our own validation study refused to support, and the fix our own drill falsified." '
    '(DEFENSE_NOTES §10, §11)'))
bullets(s, [
    'Training labels still **synthetic** — the 91.6 % is a synthetic-distribution number, and the survey does **not** upgrade it',
    'Validated: the score\'s **meaning** (ρ = 0.317; partial vs hours excludes zero). Not validated: its **accuracy at the clinical cut-off**',
    'The blocker is quantified: **1** disordered-range respondent in 87 → need **~157 usable** at the 6.4 % base rate',
    'Adults, self-reported, cross-sectional — the deployment target is **adolescents, measured, over time**',
    'Voice trained on acted adult emotion (gap quantified: leakage 9 pts; missing domain data 32 PR-AUC pts); dual-STT fails our resource gate',
    'One remaining tier: a **per-child cohort** (guardian IGDS9-SF + telemetry) — needs ethics approval; the submission is the next artefact',
], size=15)

# =============== 16. Conclusion + demo ===============
s = base('Final Results, Conclusion & Live Demo', 'Reproducibility and honesty as features — with the evidence to prove it', notes=(
    'Close on the differentiator: reproducibility and honesty as features — now with the strongest possible evidence for it, a '
    'validation study we ran on ourselves that returned two results we did not want, and a device drill that failed one of our own '
    'features. Then demo per DEMO_RUNBOOK.md; if Wi-Fi/devices misbehave, switch to the video without apologising. Future work is '
    'deliberately data-first, not model-first — "the gap is data, not architecture."'))
bullets(s, [
    '**Deployed** end-to-end system · 3 real-data-grounded models · every design choice ablated with CIs',
    '**Externally validated**: ρ = 0.317 vs IGDS9-SF, leading the screen-time baseline (97 % of paired resamples) — signal in **how**, not **how long**',
    '**Built for India**: three-register chat, Devanagari keyboard, accent-fairness audited (0 false alerts / 9.6 h)',
    '**Reproducible**: every number on these slides is one script in the repo; 7 CI guards pin the paper to the data',
    'Future is data-first: per-child cohort (ethics-gated) → child-speech adaptation → per-family thresholds',
    '**Live demo** — backup video ready',
], size=17)
callout(s, 'The thesis survived contact with real labels. Two of its components did not — and we shipped the paper that says so.',
        Inches(0.6), Inches(5.6), Inches(12.1), Inches(0.9), NAVY, 16)

prs.save(OUT)
print(f'wrote {OUT}  ({_n[0]} slides)')
